#!/usr/bin/env python
"""Build a sectioned Chinese group-meeting deck onto 汇报模板.pptx.

Reads a deck spec (JSON) and renders it with the template's own furniture --
footer and title icon -- deep-copied per slide so the result keeps the
template's look. Sections are announced by a 目录 divider before each one
(content.pptx), not by a nav bar on every slide.

    build --spec deck.json --template 汇报模板.pptx --out deck.pptx

Every content slide must carry real analysis, not just a picture: the builder
rejects any slide whose only content is a figure (see check_substance).
"""
from __future__ import annotations

import argparse
import copy
import json
import re
import sys
from datetime import date
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.packuri import PackURI
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

sys.path.insert(0, str(Path(__file__).resolve().parent))
from omml import LatexError, math_element, to_plain   # noqa: E402

GREEN = RGBColor(0x00, 0x7A, 0x49)
GREEN_D = RGBColor(0x00, 0x5C, 0x37)
INK = RGBColor(0x33, 0x33, 0x33)
BLACK = RGBColor(0x00, 0x00, 0x00)
GREY = RGBColor(0x8C, 0x8C, 0x8C)
GREY_L = RGBColor(0xBF, 0xBF, 0xBF)
RED = RGBColor(0xFF, 0x00, 0x00)      # 汇报模板2 marks figures and best results in red
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xF2, 0xF6, 0xF4)
FONT = '微软雅黑'

SECTIONS = ['文献信息', '研究背景', '本文贡献', '研究方法', '实验结果', '结论与不足']

# The 目录 lists the talk's arc, not every section. 文献信息 is front matter --
# it is on screen before the talk starts, so announcing it is noise. 研究背景 and
# 本文贡献 are one beat of the argument ("here is the gap, here is what we did
# about it") and share an entry; they still get their own slides and titles.
TOC_GROUPS = [
    ('研究背景与贡献', ('研究背景', '本文贡献')),
    ('研究方法', ('研究方法',)),
    ('实验结果', ('实验结果',)),
    ('结论与不足', ('结论与不足',)),
]
TOC_SKIP = ('文献信息',)

COVER_SLIDE, BANNER_SLIDE, THANKS_SLIDE = 0, 1, 20

# Sections are announced by a 目录 divider before each one (content.pptx), not by
# a nav bar on every slide -- so the top of every content slide is ours to use.
TITLE_Y = 0.28        # was 0.80, under the nav bar
BODY_TOP = 0.95       # was 1.45
BODY_BOTTOM = 6.9     # the footer owns everything below 7.0
ICON_SRC_Y = 0.79     # where the title icon sits in the template

TOC_NUM_X, TOC_LABEL_X = 4.88, 5.92
TOC_ROW_H, TOC_MAX_GAP = 0.79, 1.2035
TOC_OFF = RGBColor(0xD9, 0xD9, 0xD9)   # bg1 lumMod 85%, the greyed-out item

# The template ships with its author's identity on the cover and the closing
# slide. None of it belongs to the new deck, so it is replaced with a
# placeholder unless the spec supplies a real value.
PLACEHOLDER = {
    'presenter': '汇报人姓名',
    'homepage': 'https://个人主页地址',
    'school': '学校 / 学院名称',
    'logo': '学校 Logo\n（替换为本校图片）',
    'qr': '个人主页二维码\n（替换为本人二维码）',
    'paper_qr': '论文二维码\n（替换为本文链接二维码）',
}

DATE_RE = re.compile(r'\d{4}\s*年\s*\d{1,2}\s*月\s*\d{1,2}\s*日')

BOLD_RE = re.compile(r'\*\*(.+?)\*\*')
RED_RE = re.compile(r'!!(.+?)!!')
MATH_RE = re.compile(r'\$([^$]+)\$')
SEG_RE = re.compile(r'\*\*(.+?)\*\*|!!(.+?)!!|\$([^$]+)\$')


def fail(msg: str) -> None:
    raise SystemExit(f"deck spec error: {msg}")


def toc_group(section: str) -> str | None:
    """The 目录 entry a section belongs to, or None if it gets no divider."""
    for label, members in TOC_GROUPS:
        if section in members:
            return label
    return None


def check_toc_coverage() -> None:
    """Every section must be in a 目录 group or explicitly skipped -- otherwise a
    new section would quietly never announce itself."""
    covered = set(TOC_SKIP) | {s for _, members in TOC_GROUPS for s in members}
    missing = [s for s in SECTIONS if s not in covered]
    if missing:
        fail(f'sections in no 目录 group and not skipped: {missing}')


# --------------------------------------------------------------------------- #
# template furniture
# --------------------------------------------------------------------------- #

def shift_y(el, dy_emu: int) -> None:
    """Move a cloned shape vertically. For a group this is its own a:off -- the
    first a:xfrm in document order -- and never the child offsets."""
    xfrm = el.find('.//' + qn('a:xfrm'))
    off = xfrm.find(qn('a:off'))
    off.set('y', str(int(off.get('y')) + dy_emu))


def harvest_furniture(prs) -> dict[str, list]:
    """Snapshot the repeating furniture before the template body is dropped.

    The nav bar is deliberately *not* harvested: sections are announced by a
    目录 divider now. The icon is kept but moves up into the space the nav
    freed. The 目录 page takes the footer only -- its own 目录 block sits where
    the title icon would be.
    """
    store = {'footer': [], 'icon': []}
    # the footer is on every slide but the title icon only on the titled ones,
    # so take each from wherever it first turns up rather than from one slide
    for slide in list(prs.slides)[2:THANKS_SLIDE]:
        for shape in slide.shapes:
            top, left = shape.top or 0, shape.left or 0
            if shape._element.findall('.//' + qn('p:pic')):
                continue      # a picture's r:embed would not resolve on a new slide
            if not store['footer'] and top >= Inches(6.9) and left < 0:
                store['footer'].append(copy.deepcopy(shape._element))
            elif (not store['icon']
                  and abs(top - Inches(ICON_SRC_Y)) < Inches(0.06)
                  and abs(left - Inches(0.22)) < Inches(0.06)):
                el = copy.deepcopy(shape._element)
                shift_y(el, Inches(TITLE_Y) - Inches(ICON_SRC_Y))
                store['icon'].append(el)
        if store['footer'] and store['icon']:
            return store
    fail(f"template furniture not found (footer={bool(store['footer'])}, "
         f"icon={bool(store['icon'])})")


def blank_slide(prs, furniture: dict[str, list], with_icon: bool = True):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    for ph in list(slide.placeholders):
        ph._element.getparent().remove(ph._element)
    for el in furniture['footer'] + (furniture['icon'] if with_icon else []):
        slide.shapes._spTree.append(copy.deepcopy(el))
    return slide


def drop_slide(prs, index: int) -> None:
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    prs.part.drop_rel(slides[index].rId)
    xml_slides.remove(slides[index])


def move_slide(prs, old: int, new: int) -> None:
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old])
    xml_slides.insert(new, slides[old])


# --------------------------------------------------------------------------- #
# text helpers
# --------------------------------------------------------------------------- #

def style_run(run, size, bold=False, color=INK) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = FONT
    run.font.color.rgb = color
    rPr = run._r.get_or_add_rPr()
    for tag in ('latin', 'ea', 'cs'):
        el = rPr.find(qn(f'a:{tag}'))
        if el is None:
            el = rPr.makeelement(qn(f'a:{tag}'), {})
            rPr.append(el)
        el.set('typeface', FONT)
    # PowerPoint only applies CJK line-breaking rules to runs tagged Chinese;
    # without this a lone "。" or "，" can start a line
    rPr.set('lang', 'zh-CN')
    rPr.set('altLang', 'en-US')


def emit(paragraph, text: str, size: float, color=INK, bold_all=False,
         accent=GREEN) -> None:
    """Write text into a paragraph.

    ``**...**`` renders bold accent-coloured, ``!!...!!`` bold red (the best
    result in a comparison), and ``$...$`` becomes a real PowerPoint equation
    rather than literal characters.
    """
    def plain_run(s, bold, col):
        r = paragraph.add_run()
        r.text = s
        style_run(r, size, bold, col)

    pos = 0
    for m in SEG_RE.finditer(text):
        if m.start() > pos:
            plain_run(text[pos:m.start()], bold_all, color)
        if m.group(1) is not None:
            plain_run(m.group(1), True, accent)
        elif m.group(2) is not None:
            plain_run(m.group(2), True, RED)
        else:
            try:
                paragraph._p.append(math_element(m.group(3), size))
            except LatexError as exc:
                fail(f'{exc} -- in {text!r}')
        pos = m.end()
    if pos < len(text):
        plain_run(text[pos:], bold_all, color)
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('eaLnBrk', '1')
    pPr.set('hangingPunct', '1')
    pPr.set('latinLnBrk', '0')


def textbox(slide, x, y, w, h, lines, size=20, color=INK, bold_all=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.25, space_after=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        p.space_after = Pt(space_after if i < len(lines) - 1 else 0)
        emit(p, line, size, color, bold_all)
    return box


def slide_title(slide, text: str) -> None:
    textbox(slide, 0.88, TITLE_Y, 8.0, 0.55, [text], size=24, color=GREEN, bold_all=True)


def rounded(slide, x, y, w, h, fill=WASH, line=GREEN, radius=0.06):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1.25)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    sh.text_frame.text = ''
    return sh


def fit_picture(slide, path, x, y, w, h):
    """Place an image inside a box, preserving aspect, centred."""
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    return slide.shapes.add_picture(str(path), Inches(x + (w - pw) / 2),
                                    Inches(y + (h - ph) / 2), Inches(pw), Inches(ph))


def as_rendered(text: str) -> str:
    """What the audience actually sees: markers gone, math laid out flat."""
    text = RED_RE.sub(r'\1', BOLD_RE.sub(r'\1', text))
    return MATH_RE.sub(lambda m: to_plain(m.group(1)), text)


def text_width(text: str, size: float) -> float:
    """Rendered width in inches. A CJK glyph is ~1em, latin ~0.55em."""
    em = size / 72
    return sum(em if ord(c) > 0x2E80 else em * 0.55 for c in as_rendered(text))


def wrapped_height(lines, text_w_in, size, spacing=1.25, space_after=7) -> float:
    """Estimate rendered height. A CJK glyph is ~1em wide, so characters per line
    follow from the box width; markers and LaTeX syntax do not render."""
    cpl = max(4, int(text_w_in * 72 / size))
    rows = sum(max(1, -(-len(as_rendered(l)) // cpl)) for l in lines)
    return rows * size * spacing / 72 + (len(lines) - 1) * space_after / 72


def fit_size(lines, w_in, h_in, size, spacing=1.35, space_after=12, floor=11) -> float:
    """Shrink until the text fits its box. A fixed size plus a list that grew by
    one item is how a panel silently spills past the footer."""
    while size > floor and wrapped_height(lines, w_in, size, spacing,
                                          space_after) > h_in:
        size -= 0.5
    return size


def bullet_box(slide, x, y, w, bullets, size=20) -> float:
    """The template's 要点框: a light band of green-dotted bullets. Returns bottom."""
    lines = ['• ' + b for b in bullets]
    h = wrapped_height(lines, w - 0.4, size) + 0.34
    rounded(slide, x, y, w, h, fill=WASH, line=None)
    textbox(slide, x + 0.18, y + 0.16, w - 0.36, h - 0.32, lines,
            size=size, spacing=1.2, space_after=7)
    return y + h


CALLOUT_HEAD_SIZE = 18


def callout_metrics(w, text, size, head):
    """Tab width, body width and height of a callout -- one source of truth, so
    a caller that reserves room for one gets the same answer the builder uses."""
    lines = text if isinstance(text, list) else [text]
    tab_w = max(1.35, chevron_width(head, CALLOUT_HEAD_SIZE))
    body_w = w - tab_w - 0.35
    h = max(0.9, wrapped_height(lines, body_w, size, space_after=5) + 0.3)
    return tab_w, body_w, h


def analysis_callout(slide, x, y, w, text, size=20, head='结果分析') -> float:
    """The template's chevron callout: green tab + outlined body.

    The tab is sized from the head so it never wraps, and the body starts where
    the tab actually ends rather than at a fixed offset.
    """
    lines = text if isinstance(text, list) else [text]
    tab_w, body_w, h = callout_metrics(w, text, size, head)
    rounded(slide, x, y, w, h, fill=WHITE, line=GREEN)
    chevron(slide, x, y, tab_w, h, head, size=CALLOUT_HEAD_SIZE)
    textbox(slide, x + tab_w + 0.15, y + 0.1, body_w, h - 0.2, lines,
            size=size, anchor=MSO_ANCHOR.MIDDLE, spacing=1.2, space_after=5)
    return y + h


def note_height(text, w, size=15) -> float:
    """How much room a caption really needs -- a long one wraps to two lines and
    would otherwise be drawn straight through the footer."""
    return max(0.42, wrapped_height([text], w - 0.12, size, spacing=1.2,
                                    space_after=0) + 0.1)


def figure_note(slide, x, y, w, text, size=15):
    """图注 defaults to black -- it is caption text, not an accent."""
    textbox(slide, x, y, w, note_height(text, w, size), [text], size=size,
            color=BLACK, align=PP_ALIGN.CENTER, spacing=1.2, space_after=0)


def annotate(slide, box, annots, size=13):
    """Red notes pinned onto a figure, as 汇报模板2 does.

    `box` is the picture's placed rect; each annot's `at` is a fraction of it,
    so a note tracks its panel no matter how the figure was scaled.
    """
    x, y, w, h = box
    for a in annots:
        fx, fy = a.get('at', (0.5, 0.0))
        aw = a.get('w', 2.3)
        lines = a['text'] if isinstance(a['text'], list) else [a['text']]
        ah = wrapped_height(lines, aw - 0.12, size, spacing=1.1, space_after=2) + 0.1
        ax = min(max(x + fx * w - aw / 2, 0.12), 13.2 - aw)
        ay = min(max(y + fy * h - ah / 2, BODY_TOP - 0.15), BODY_BOTTOM - ah)
        textbox(slide, ax, ay, aw, ah, lines, size=size, color=RED, bold_all=True,
                align=PP_ALIGN.CENTER, spacing=1.1, space_after=2)


def place_figure(slide, spec_fig, x, y, w, h):
    """A figure plus its optional red annotations. Accepts a path or a dict."""
    if isinstance(spec_fig, dict):
        path, annots = spec_fig['path'], spec_fig.get('annotations', [])
    else:
        path, annots = spec_fig, []
    pic = fit_picture(slide, path, x, y, w, h)
    box = (pic.left / 914400, pic.top / 914400,
           pic.width / 914400, pic.height / 914400)
    if annots:
        annotate(slide, box, annots)
    return box


# A head is a label, not prose: 「结果分析」 broken across two lines as 「结果分 /
# 析」 looks like a bug. Every head shape is sized from its text and has wrapping
# switched off, so it always renders on one line.
TAB_PAD = 0.34      # padding either side of a pill's text
TIP_PAD = 0.78      # padding for a PENTAGON: the arrow tip eats the right side


def head_text(shape, text, size, colour) -> None:
    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    emit(p, text, size, colour, True)


def pill_width(text, size=18) -> float:
    return text_width(text, size) + TAB_PAD * 2


def pill(slide, x, y, w, h, text, fill=GREEN, size=18):
    """The rounded green label 汇报模板2 puts at the head of an analysis panel."""
    sh = rounded(slide, x, y, max(w, pill_width(text, size)), h,
                 fill=fill, line=None, radius=0.3)
    head_text(sh, text, size, WHITE)
    return sh


def chevron_width(text, size=17) -> float:
    return text_width(text, size) + TIP_PAD + TAB_PAD


def chevron(slide, x, y, w, h, text, size=17):
    w = max(w, chevron_width(text, size))
    tab = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    tab.fill.solid()
    tab.fill.fore_color.rgb = WHITE
    tab.line.color.rgb = GREEN
    tab.line.width = Pt(1.25)
    tab.shadow.inherit = False
    # the tip is on the right, so nudge the text left of it to stay centred
    tf = tab.text_frame
    head_text(tab, text, size, GREEN)
    tf.margin_right = Inches(TIP_PAD / 2)
    return tab


def arrow(slide, x, y, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = GREEN
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def statement(slide, x, y, w, lines, size=19) -> float:
    """The full-width claim banner that opens 汇报模板2's summary slides."""
    lines = lines if isinstance(lines, list) else [lines]
    h = wrapped_height(lines, w - 0.7, size, spacing=1.3, space_after=6) + 0.36
    rounded(slide, x, y, w, h, fill=WHITE, line=GREEN, radius=0.05)
    textbox(slide, x + 0.3, y + 0.14, w - 0.6, h - 0.28, lines, size=size,
            color=GREEN_D, bold_all=True, anchor=MSO_ANCHOR.MIDDLE,
            spacing=1.3, space_after=6)
    return y + h


# --------------------------------------------------------------------------- #
# slide builders
# --------------------------------------------------------------------------- #

def retext(shape, text: str) -> bool:
    """Rewrite a shape's (or table cell's) first non-empty paragraph, keeping
    its formatting.

    The template's banner keeps its authors in paragraph 8 of a full-slide frame,
    so paragraph 0 is not a safe assumption.
    """
    if not getattr(shape, 'has_text_frame', True):
        return False
    for para in shape.text_frame.paragraphs:
        runs = para.runs
        if not runs:
            continue
        runs[0].text = text
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
        return True
    return False


def placeholder_box(slide, x, y, w, h, label: str) -> None:
    """A dashed box standing in for an asset the presenter must supply."""
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = WASH
    sh.line.color.rgb = GREY
    sh.line.width = Pt(1.0)
    sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(label.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        emit(p, line, 12 if i else 14, GREY, i == 0)


def replace_asset(slide, pic, path, label: str) -> None:
    """Swap a template image for the user's, or for a labelled placeholder.

    The template's logo and QR codes encode its author's identity -- the QR in
    particular still resolves to their links, which no amount of retexting
    fixes -- so an unsupplied asset must not simply be left in place.
    """
    x, y = pic.left / 914400, pic.top / 914400
    w, h = pic.width / 914400, pic.height / 914400
    pic._element.getparent().remove(pic._element)
    if path:
        fit_picture(slide, path, x, y, w, h)
    else:
        placeholder_box(slide, x, y, w, h, label)


def build_identity(prs, meta, index: int) -> None:
    """Fill the presenter's own details on the cover / closing slide.

    Defaults are placeholders on purpose: shipping the template author's name,
    school and homepage in someone else's deck is worse than an obvious blank.
    """
    slide = prs.slides[index]
    presenter = meta.get('presenter') or PLACEHOLDER['presenter']
    homepage = meta.get('homepage') or PLACEHOLDER['homepage']
    school = meta.get('school') or PLACEHOLDER['school']
    for sh in list(slide.shapes):
        top, left = (sh.top or 0) / 914400, (sh.left or 0) / 914400
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if left < 1.0 and top < 1.0:
                replace_asset(slide, sh, meta.get('logo'), PLACEHOLDER['logo'])
            elif left > 10.0 and top < 1.0:
                replace_asset(slide, sh, meta.get('qr'), PLACEHOLDER['qr'])
        elif sh.has_table:
            cells = list(sh.table.rows[0].cells)   # 汇报人 | ： | 名字
            if '汇报人' in cells[0].text:
                retext(cells[-1], presenter)
        elif sh.has_text_frame:
            text = sh.text_frame.text.strip()
            if re.match(r'https?://', text):
                retext(sh, homepage)
            elif top > 4.0 and ('学院' in text or '大学' in text):
                retext(sh, school)


def build_cover(prs, meta) -> None:
    slide = prs.slides[COVER_SLIDE]
    for sh in slide.shapes:
        if (sh.has_text_frame and sh.text_frame.text.strip()
                and sh.width > Inches(6) and Inches(2.0) < sh.top < Inches(3.0)):
            if retext(sh, meta['topic_zh']):
                break
    else:
        print('warning: cover title not found', file=sys.stderr)
    build_identity(prs, meta, COVER_SLIDE)


def build_banner(prs, meta) -> None:
    """Fill the paper banner. Targets are picked by role, not by position alone:
    the venue box and the outer frame share the same top edge."""
    slide = prs.slides[BANNER_SLIDE]
    todo = {'title_en': meta.get('title_en'), 'authors': meta.get('authors'),
            'venue_line': meta.get('venue_line'), 'url': meta.get('url')}
    done = {k: False for k in todo}
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            replace_asset(slide, sh, meta.get('paper_qr'), PLACEHOLDER['paper_qr'])
            continue
        if not sh.has_text_frame or not sh.text_frame.text.strip():
            continue
        top, w = sh.top or 0, sh.width or 0
        if Inches(2.0) < top < Inches(3.4):
            key = 'title_en'
        elif top < Inches(0.75) and w > Inches(9):     # full-slide frame: authors
            key = 'authors'
        elif top < Inches(0.75):                        # narrow box: venue
            key = 'venue_line'
        elif top > Inches(5.8):
            key = 'url'
        else:
            continue
        if todo.get(key) and not done[key]:
            done[key] = retext(sh, todo[key])
    missed = [k for k, v in todo.items() if v and not done[k]]
    if missed:
        print(f'warning: banner fields not filled: {missed}', file=sys.stderr)


def build_info(prs, spec, meta, furn) -> None:
    slide = blank_slide(prs, furn)
    slide_title(slide, spec.get('title', '文献信息'))
    cover = meta.get('cover_page') or meta.get('cover_crop')
    if cover:
        fit_picture(slide, cover, 0.30, BODY_TOP, 5.9, 5.65)
        figure_note(slide, 0.30, BODY_TOP + 5.7, 5.9, spec.get('cover_note', '论文首页'))
    rows = list(meta.get('info', {}).items())
    y = BODY_TOP + 0.07
    rounded(slide, 6.55, BODY_TOP, 6.5, 0.42 + 0.62 * len(rows), fill=WASH, line=None)
    for k, v in rows:
        textbox(slide, 6.75, y, 2.3, 0.5, [f'【{k}】'], size=20, color=GREEN, bold_all=True)
        textbox(slide, 9.05, y, 3.9, 0.5, [f'{v}'], size=20, bold_all=True)
        y += 0.62
    if spec.get('bullets'):
        bullet_box(slide, 6.55, y + 0.25, 6.5, spec['bullets'], size=18)


def build_bullets(prs, spec, meta, furn) -> None:
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    bottom = bullet_box(slide, 0.35, BODY_TOP, 12.6, spec['bullets'])
    figs = spec.get('figures', [])
    notes = spec.get('notes', [])
    if not figs:
        if spec.get('analysis'):
            analysis_callout(slide, 0.35, bottom + 0.35, 12.6, spec['analysis'],
                             head=spec.get('analysis_head', '要点'))
        return
    top = bottom + 0.3
    gap, n = 0.25, len(figs)
    w = (12.6 - gap * (n - 1)) / n
    nh = max((note_height(t, w) for t in notes), default=0)
    avail_h = BODY_BOTTOM - top - nh
    for i, f in enumerate(figs):
        x = 0.35 + i * (w + gap)
        place_figure(slide, f, x, top, w, avail_h)
        if i < len(notes):
            figure_note(slide, x, top + avail_h + 0.04, w, notes[i])


def build_cards(prs, spec, meta, furn) -> None:
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    cards = spec['cards']
    n = len(cards)
    gap = 0.35
    w = (12.6 - gap * (n - 1)) / n
    y = BODY_TOP + 0.3
    # the callout has to fit above the footer, so the cards give up the room --
    # exactly as much as it will actually take, not a guess
    ch = (callout_metrics(12.6, spec['analysis'], 18,
                          spec.get('analysis_head', '小结'))[2] + 0.2
          if spec.get('analysis') else 0)
    h = 6.9 - y - ch
    for i, c in enumerate(cards):
        x = 0.35 + i * (w + gap)
        rounded(slide, x, y, w, h, fill=WHITE, line=GREEN, radius=0.08)
        lines = [f"**{c['lead']}**：{c['text']}"] if c.get('lead') else [c['text']]
        if c.get('points'):
            lines += ['• ' + p for p in c['points']]
        textbox(slide, x + 0.28, y + 0.3, w - 0.56, h - 0.6, lines,
                size=c.get('size', 18), anchor=MSO_ANCHOR.MIDDLE, spacing=1.35, space_after=9)
    if spec.get('analysis'):
        analysis_callout(slide, 0.35, y + h + 0.2, 12.6, spec['analysis'],
                         head=spec.get('analysis_head', '小结'), size=18)


def build_figure_analysis(prs, spec, meta, furn) -> None:
    """Figure on one side, bullets + analysis on the other."""
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    figs = spec.get('figures', [])
    side = spec.get('figure_side', 'right')
    fw = spec.get('figure_width', 7.3)
    tw = 12.98 - fw - 0.7
    fx = 13.0 - fw - 0.35 if side == 'right' else 0.35
    tx = 0.35 if side == 'right' else fx + fw + 0.35
    notes = spec.get('notes', [])
    top, bottom, gap = BODY_TOP, BODY_BOTTOM, 0.2
    n = max(1, len(figs))
    each = (bottom - top - gap * (n - 1)) / n
    for i, fig in enumerate(figs):
        yy = top + i * (each + gap)
        nh = note_height(notes[i], fw) if i < len(notes) else 0
        img_h = each - nh                            # the note lives in the slot,
        place_figure(slide, fig, fx, yy, fw, img_h)  # not in the gap to the next
        if nh:
            figure_note(slide, fx, yy + img_h + 0.02, fw, notes[i])
    y = top
    if spec.get('bullets'):
        y = bullet_box(slide, tx, y, tw, spec['bullets'], size=spec.get('size', 18)) + 0.3
    if spec.get('analysis'):
        analysis_callout(slide, tx, y, tw, spec['analysis'],
                         head=spec.get('analysis_head', '结果分析'),
                         size=spec.get('size', 18))


def build_numbered(prs, spec, meta, furn) -> None:
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    items = spec['items']
    y0, h = BODY_TOP, BODY_BOTTOM - BODY_TOP - 0.1
    rounded(slide, 0.45, y0, 12.45, h, fill=WHITE, line=GREEN, radius=0.05)
    each = h / len(items)
    for i, it in enumerate(items):
        y = y0 + i * each
        textbox(slide, 0.85, y + each / 2 - 0.55, 1.0, 1.1, [str(i + 1)],
                size=54, color=GREEN, bold_all=True, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, 1.95, y + 0.22, 10.6, 0.5, [it['title']], size=21,
                color=GREEN_D, bold_all=True)
        textbox(slide, 1.95, y + 0.78, 10.6, each - 1.0, it['desc'] if isinstance(it['desc'], list) else [it['desc']],
                size=18, spacing=1.3)


def harvest_toc_block(path):
    """The 目录 badge from content.pptx: green squares plus the white 目录.

    It is a group of autoshapes with no picture in it, so it deep-copies into
    another deck cleanly.
    """
    if not Path(path).exists():
        print(f'warning: {path} not found; 目录 pages will have no badge',
              file=sys.stderr)
        return None
    src = Presentation(path)
    for shape in src.slides[0].shapes:
        if (shape.shape_type == MSO_SHAPE_TYPE.GROUP
                and (shape.left or 0) < Inches(1.2) and (shape.top or 0) < Inches(1.2)):
            if shape._element.findall('.//' + qn('p:pic')):
                break
            return copy.deepcopy(shape._element)
    print(f'warning: no 目录 badge found in {path}', file=sys.stderr)
    return None


def build_toc(prs, entries, active, furn, badge) -> None:
    """The 目录 divider that opens each part of the talk (content.pptx).

    Every entry is listed every time; only the one being entered is black. This
    replaces the per-slide nav bar -- the audience gets a real pause at each
    boundary instead of a strip they stop reading after slide three.
    """
    slide = blank_slide(prs, furn, with_icon=False)   # the badge owns that corner
    if badge is not None:
        slide.shapes._spTree.append(copy.deepcopy(badge))
    n = len(entries)
    gap = min(TOC_MAX_GAP, 4.9 / max(1, n - 1))
    y = (7.5 - ((n - 1) * gap + TOC_ROW_H)) / 2
    for i, sec in enumerate(entries, 1):
        num = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(TOC_NUM_X),
                                     Inches(y), Inches(TOC_ROW_H), Inches(TOC_ROW_H))
        num.fill.solid()
        num.fill.fore_color.rgb = GREEN
        num.line.fill.background()
        num.shadow.inherit = False
        head_text(num, f'{i:02d}', 28, WHITE)
        box = textbox(slide, TOC_LABEL_X, y + 0.145, 6.5, 0.5, [sec], size=24,
                      color=BLACK if sec == active else TOC_OFF, bold_all=True,
                      anchor=MSO_ANCHOR.MIDDLE)
        for para in box.text_frame.paragraphs:
            for run in para.runs:
                run._r.get_or_add_rPr().set('spc', '300')   # the template's tracking
        y += gap


def build_process(prs, spec, meta, furn) -> None:
    """图 + 底部流程条: the pipeline as named steps, not one summary sentence."""
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    steps = spec['steps']
    strip_h = 1.15
    top, bottom = BODY_TOP, 6.9 - strip_h - 0.25
    figs = spec.get('figures', [])
    notes = spec.get('notes', [])
    n = max(1, len(figs))
    gap = 0.25
    w = (12.6 - gap * (n - 1)) / n
    for i, fig in enumerate(figs):
        x = 0.35 + i * (w + gap)
        nh = note_height(notes[i], w) if i < len(notes) else 0
        fh = bottom - top - nh
        place_figure(slide, fig, x, top, w, fh)
        if nh:
            figure_note(slide, x, top + fh + 0.02, w, notes[i])
    y = 6.9 - strip_h
    head = spec.get('steps_head', '流程')
    tab_w = max(1.5, chevron_width(head))
    rounded(slide, 0.35, y, 12.6, strip_h, fill=WHITE, line=GREEN, radius=0.05)
    chevron(slide, 0.35, y, tab_w, strip_h, head)
    # steps and the arrows between them share the strip's remaining width
    n = len(steps)
    aw = 0.5
    each = (12.6 - tab_w - 0.25 - (n - 1) * aw) / n
    x = 0.35 + tab_w + 0.25
    for i, s in enumerate(steps):
        textbox(slide, x, y + 0.08, each, strip_h - 0.16,
                s if isinstance(s, list) else [s], size=spec.get('size', 16),
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, spacing=1.15)
        x += each
        if i < n - 1:
            arrow(slide, x + 0.06, y + strip_h / 2 - 0.14, aw - 0.12, 0.28)
            x += aw


def build_tree(prs, spec, meta, furn) -> None:
    """一句论断 → 箭头 → 并列分支: why the existing options each fall short."""
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    y = statement(slide, 0.35, BODY_TOP, 12.6, spec['statement'])
    branches = spec['branches']
    n = len(branches)
    gap = 0.4
    w = (12.6 - gap * (n - 1)) / n
    stem, by = y + 0.15, y + 0.75
    h = 6.9 - by
    for i, b in enumerate(branches):
        x = 0.35 + i * (w + gap)
        cx = x + w / 2
        conn = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(cx - 0.09),
                                      Inches(stem), Inches(0.18), Inches(0.5))
        conn.fill.solid()
        conn.fill.fore_color.rgb = GREEN
        conn.line.fill.background()
        conn.shadow.inherit = False
        rounded(slide, x, by, w, h, fill=WHITE, line=GREEN, radius=0.06)
        lines = [f"**{b['lead']}**：{b['text']}"] if b.get('lead') else [b['text']]
        if b.get('points'):
            lines += ['• ' + p for p in b['points']]
        textbox(slide, x + 0.25, by + 0.22, w - 0.5, h - 0.44, lines,
                size=spec.get('size', 17), spacing=1.3, space_after=8)


def build_figure_cards(prs, spec, meta, furn) -> None:
    """一句论断 + 每张卡「图 + 小标题 + 说明」: three findings, three figures."""
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    y = BODY_TOP
    if spec.get('statement'):
        y = statement(slide, 0.35, y, 12.6, spec['statement']) + 0.25
    cards = spec['cards']
    n = len(cards)
    gap = 0.35
    w = (12.6 - gap * (n - 1)) / n
    h = 6.9 - y
    size = spec.get('size', 16)
    for i, c in enumerate(cards):
        x = 0.35 + i * (w + gap)
        rounded(slide, x, y, w, h, fill=WHITE, line=GREEN, radius=0.06)
        desc = c['desc'] if isinstance(c['desc'], list) else [c['desc']]
        dh = wrapped_height(desc, w - 0.7, size, spacing=1.3) + 0.1
        th = 0.42
        place_figure(slide, c['figure'], x + 0.18, y + 0.18,
                     w - 0.36, h - 0.36 - th - dh)
        ty = y + h - 0.18 - dh - th
        textbox(slide, x + 0.25, ty, w - 0.5, th, [c['title']], size=19,
                color=GREEN_D, bold_all=True, align=PP_ALIGN.CENTER)
        textbox(slide, x + 0.35, ty + th, w - 0.7, dh, desc, size=size, spacing=1.3)


def build_panel(prs, spec, meta, furn) -> None:
    """大图铺开 + 右侧编号分析面板: room for several readings of one result."""
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    pw = spec.get('panel_width', 4.3)
    px = 13.0 - pw - 0.35
    figs = spec.get('figures', [])
    notes = spec.get('notes', [])
    top, bottom = BODY_TOP, BODY_BOTTOM
    fw_all = px - 0.7
    n = max(1, len(figs))
    gap = 0.2
    w = (fw_all - gap * (n - 1)) / n
    for i, fig in enumerate(figs):
        x = 0.35 + i * (w + gap)
        nh = note_height(notes[i], w) if i < len(notes) else 0
        fh = bottom - top - nh
        place_figure(slide, fig, x, top, w, fh)
        if nh:
            figure_note(slide, x, top + fh + 0.02, w, notes[i])
    rounded(slide, px, top, pw, bottom - top, fill=WHITE, line=GREEN, radius=0.05)
    pill(slide, px + 0.35, top + 0.25, pw - 0.7, 0.52,
         spec.get('panel_head', '结果分析'))
    lines = [f'{i}. {t}' for i, t in enumerate(spec['items'], 1)]
    avail = bottom - top - 1.25
    size = fit_size(lines, pw - 0.7, avail, spec.get('size', 17))
    textbox(slide, px + 0.35, top + 1.0, pw - 0.7, avail, lines,
            size=size, spacing=1.35, space_after=12)


def build_proscons(prs, spec, meta, furn) -> None:
    """创新点 / 局限性 as two banded lists -- the template's closing contrast."""
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    y = BODY_TOP
    if spec.get('statement'):
        y = statement(slide, 0.35, y, 12.6, spec['statement']) + 0.3
    bands = [(spec.get('pros_head', '创新点'), spec['pros'], GREEN),
             (spec.get('cons_head', '局限性'), spec['cons'], GREY_L)]
    # layout: pros list / pros band / gap / cons band / cons list. The pill is
    # taller than its band, so the two bands need clearance or the pills collide.
    band_h, pill_h, gap = 0.42, 0.62, 0.75
    list_h = (6.9 - y - gap - band_h) / 2
    # the lists clear the widest pill, so a long head cannot run into them
    tab_w = max(2.2, max(pill_width(h) for h, _, _ in bands))
    list_x = 0.45 + tab_w + 0.3
    list_w = 12.95 - list_x - 0.35
    # both lists share one size so the pair reads as a matched contrast
    size = min(fit_size([f'{i}. {t}' for i, t in enumerate(items, 1)],
                        list_w, list_h, spec.get('size', 17), space_after=9)
               for _, items, _ in bands)

    def band(by, head, colour):
        rounded(slide, 0.35, by, 12.6, band_h, fill=colour, line=None, radius=0.2)
        pill(slide, 0.45, by - (pill_h - band_h) / 2, tab_w, pill_h, head, fill=colour)

    def listing(ly, items):
        textbox(slide, list_x, ly, list_w, list_h,
                [f'{i}. {t}' for i, t in enumerate(items, 1)], size=size,
                anchor=MSO_ANCHOR.MIDDLE, spacing=1.35, space_after=9)

    listing(y, bands[0][1])
    band(y + list_h, bands[0][0], bands[0][2])
    cy = y + list_h + gap
    band(cy, bands[1][0], bands[1][2])
    listing(cy + band_h, bands[1][1])


def build_dual(prs, spec, meta, furn) -> None:
    """两图并列，各配各的分析: two results that answer different questions."""
    slide = blank_slide(prs, furn)
    slide_title(slide, spec['title'])
    panels = spec['panels']
    n = len(panels)
    gap = 0.4
    w = (12.6 - gap * (n - 1)) / n
    size = spec.get('size', 16)
    top = BODY_TOP
    # every panel's callout starts at the same y, so the row reads as a row
    ch = max(callout_metrics(w, p['analysis'], size, p.get('head', '分析结果'))[2]
             for p in panels)
    cy = 6.9 - ch
    for i, p in enumerate(panels):
        x = 0.35 + i * (w + gap)
        nh = note_height(p['note'], w) if p.get('note') else 0
        fh = cy - 0.25 - top - nh
        place_figure(slide, p['figure'], x, top, w, fh)
        if nh:
            figure_note(slide, x, top + fh + 0.02, w, p['note'])
        analysis_callout(slide, x, cy, w, p['analysis'],
                         head=p.get('head', '分析结果'), size=size)


BUILDERS = {'info': build_info, 'bullets': build_bullets, 'cards': build_cards,
            'figure_analysis': build_figure_analysis, 'numbered': build_numbered,
            'process': build_process, 'tree': build_tree,
            'figure_cards': build_figure_cards, 'panel': build_panel,
            'proscons': build_proscons, 'dual': build_dual}

# a slide type carries its own content; these are the keys that count as
# "said something" for the substance check
TEXTY = ('bullets', 'analysis', 'steps', 'branches', 'items', 'cards',
         'panels', 'pros', 'statement')


# --------------------------------------------------------------------------- #
# validation
# --------------------------------------------------------------------------- #

REQUIRED = {
    'info': (), 'bullets': ('bullets',), 'cards': ('cards',),
    'figure_analysis': ('bullets',), 'numbered': ('items',),
    'process': ('steps', 'figures'), 'tree': ('statement', 'branches'),
    'figure_cards': ('cards',), 'panel': ('items', 'figures'),
    'proscons': ('pros', 'cons'), 'dual': ('panels',),
}


def check_substance(spec) -> None:
    """A group-meeting slide must say something, not just show a figure."""
    for i, s in enumerate(spec['slides'], 1):
        t = s.get('type')
        if t not in BUILDERS:
            fail(f"slide {i}: unknown type {t!r} (use {sorted(BUILDERS)})")
        for key in REQUIRED[t]:
            if not s.get(key):
                fail(f"slide {i}: type {t!r} needs {key!r}")
        # every slide needs a section: it decides where a 目录 divider goes
        if not s.get('section'):
            fail(f"slide {i}: missing section (one of {SECTIONS})")
        if s['section'] not in SECTIONS:
            fail(f"slide {i}: unknown section {s['section']!r}")
        if t in ('bullets', 'figure_analysis'):
            text_blocks = len(s.get('bullets', [])) + (1 if s.get('analysis') else 0)
            if s.get('figures') and text_blocks < 2:
                fail(f"slide {i} ({s.get('title')!r}): a figure slide needs bullets AND "
                     f"an analysis callout -- summarise the paper, do not just show "
                     f"the figure")
            if not s.get('figures') and not s.get('bullets'):
                fail(f"slide {i}: empty slide")


def check_variety(spec) -> None:
    """研究方法 and 实验结果 must not be the same slide N times over.

    The whole point of the second template is that a pipeline wants a 流程条, a
    comparison wants a 分析面板, and three findings want three 图卡 -- reaching
    for figure_analysis every time is the rut this check exists to break.
    """
    by_section: dict[str, list] = {}
    for i, s in enumerate(spec['slides'], 1):
        if s.get('section'):
            by_section.setdefault(s['section'], []).append((i, s['type']))
    for section in ('研究方法', '实验结果'):
        entries = by_section.get(section, [])
        types = [t for _, t in entries]
        if len(types) >= 3 and len(set(types)) < 2:
            fail(f"{section}: {len(types)} slides all of type {types[0]!r} -- vary the "
                 f"layout ({sorted(set(BUILDERS) - {'info'})})")
        run, start = 1, 0
        for j in range(1, len(types)):
            run = run + 1 if types[j] == types[j - 1] else 1
            if run == 1:
                start = j
            if run > 3:
                fail(f"{section}: slides {entries[start][0]}-{entries[j][0]} are four "
                     f"{types[j]!r} slides in a row -- break the pattern")


def check_emphasis(spec) -> None:
    """实验结果 marks the winning number in red; a table of numbers with nothing
    picked out makes the audience hunt for the point."""
    hits = 0
    for s in spec['slides']:
        if s.get('section') == '实验结果':
            hits += len(RED_RE.findall(json.dumps(s, ensure_ascii=False)))
    if not hits:
        fail("实验结果: no !!红色强调!! anywhere -- mark the best result in each "
             "comparison so the audience sees it without reading the table")


HEAD_KEYS = ('analysis_head', 'panel_head', 'steps_head', 'pros_head', 'cons_head')
MAX_HEAD_W = 2.6      # inches -- about 8 CJK characters at 18pt


def check_heads(spec) -> None:
    """A head is a label. The shapes grow to keep it on one line, so an essay in
    an `analysis_head` would squeeze out the text it is labelling."""
    def check(head, where):
        if head and chevron_width(head, CALLOUT_HEAD_SIZE) > MAX_HEAD_W:
            fail(f'{where}: head {head!r} is too long -- a few characters, like '
                 f'结果分析 / 关键设计 / 小结')
    for i, s in enumerate(spec['slides'], 1):
        for key in HEAD_KEYS:
            check(s.get(key), f'slide {i} {key}')
        for p in s.get('panels', []):
            check(p.get('head'), f'slide {i} panel head')


def check_math(spec) -> None:
    """Reject bad LaTeX up front rather than half way through a build."""
    def walk(node):
        if isinstance(node, str):
            if node.count('$') % 2:
                fail(f'unbalanced $ in {node!r}')
            for m in MATH_RE.finditer(node):
                try:
                    math_element(m.group(1))
                except LatexError as exc:
                    fail(f'{exc} -- in {node!r}')
        elif isinstance(node, dict):
            for v in node.values():
                walk(v)
        elif isinstance(node, list):
            for v in node:
                walk(v)
    walk(spec['slides'])


def retext_dates(tf, date: str) -> int:
    """Rewrite any 「YYYY年M月D日」 in a text frame, keeping its formatting."""
    hits = 0
    for para in tf.paragraphs:
        runs = para.runs
        if not runs:
            continue
        joined = ''.join(r.text for r in runs)
        if not DATE_RE.search(joined):
            continue
        runs[0].text = DATE_RE.sub(date, joined)   # the label may be split
        for r in runs[1:]:                         # across several runs
            r._r.getparent().remove(r._r)
        hits += 1
    return hits


def set_date(prs, date: str) -> int:
    """Stamp today's date over the template's own.

    The footer date is a table on every template slide -- and the footer is
    cloned as furniture -- so this has to run before harvest_furniture, and the
    template's 2026年1月21日 would otherwise ship on all 19 slides.
    """
    hits = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            if sh.has_text_frame:
                hits += retext_dates(sh.text_frame, date)
            elif getattr(sh, 'has_table', False):
                for row in sh.table.rows:
                    for cell in row.cells:
                        hits += retext_dates(cell.text_frame, date)
    if not hits:
        print('warning: no 年月日 found to restamp; the footer may show the '
              "template's date", file=sys.stderr)
    return hits


def set_total_pages(prs, total: int) -> int:
    """The master hard-codes 共N页; fix it or every slide shows the wrong count.

    PowerPoint splits that label across runs (第 / ‹#› / 页 / 共 / 22 / 页), so
    match on the joined paragraph text and patch whichever run holds the number.
    """
    hits = 0
    for master in prs.slide_masters:
        for shape in master.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                runs = para.runs
                joined = ''.join(r.text for r in runs)
                if not re.search(r'共\s*\d+\s*页', joined):
                    continue
                for i, run in enumerate(runs):
                    if re.search(r'共\d+页', run.text):          # all in one run
                        run.text = re.sub(r'共\d+页', f'共{total}页', run.text)
                        hits += 1
                    elif run.text.strip().isdigit() and i and runs[i - 1].text.endswith('共'):
                        run.text = str(total)                    # split across runs
                        hits += 1
    if not hits:
        print('warning: could not find 共N页 in the master; page count may be stale',
              file=sys.stderr)
    return hits


def main(argv=None) -> None:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument('--spec', required=True)
    ap.add_argument('--template', required=True)
    ap.add_argument('--out', required=True)
    ap.add_argument('--content', default=str(Path(__file__).resolve().parent.parent
                                             / 'assets' / 'content.pptx'),
                    help='deck supplying the 目录 badge')
    args = ap.parse_args(argv)

    spec = json.loads(Path(args.spec).read_text(encoding='utf-8'))
    check_toc_coverage()
    check_substance(spec)
    check_variety(spec)
    check_emphasis(spec)
    check_heads(spec)
    check_math(spec)

    prs = Presentation(args.template)
    meta = spec['meta']
    today = date.today()
    set_date(prs, meta.get('date') or f'{today.year}年{today.month}月{today.day}日')
    build_cover(prs, meta)
    build_banner(prs, meta)
    build_identity(prs, meta, THANKS_SLIDE)   # while the template index is still valid

    furniture = harvest_furniture(prs)   # before the body slides are dropped

    keep_front = 2                       # cover + banner
    for _ in range(len(prs.slides) - keep_front - 1):
        drop_slide(prs, keep_front)      # template body; thanks slide is last
    thanks_at = len(prs.slides) - 1
    # python-pptx names each new slide /ppt/slides/slide{count+1}.xml, so once
    # the deck passes 20 slides it hands out slide21.xml -- the partname the
    # retained 谢谢大家 slide already holds, and the .pptx gets two entries with
    # the same name. Move it out of the counter's range.
    prs.slides[thanks_at].part.partname = PackURI('/ppt/slides/slide9000.xml')

    badge = harvest_toc_block(args.content)
    used = [label for label, members in TOC_GROUPS
            if any(sl['section'] in members for sl in spec['slides'])]
    current = None
    for sl in spec['slides']:
        group = toc_group(sl['section'])
        if group and group != current:    # a 目录 page opens every 目录 entry
            build_toc(prs, used, group, furniture, badge)
        current = group                   # 文献信息 (no group) still ends the run
        BUILDERS[sl['type']](prs, sl, meta, furniture)
    # the thanks slide came from the template and must end up last
    move_slide(prs, thanks_at, len(prs.slides) - 1)

    set_total_pages(prs, len(prs.slides))
    Path(args.out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)
    print(f"Wrote {args.out}  ({len(prs.slides)} slides)")


if __name__ == '__main__':
    main()
